import streamlit as st
import google.generativeai as genai
import pandas as pd
import io
import json
import re
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt

# --- 1. CONFIG & SESSION STATE ---
st.set_page_config(page_title="Science Lab", layout="wide")

if 'scenario' not in st.session_state: st.session_state.scenario = None
if 'current_df' not in st.session_state: st.session_state.current_df = None
if 'model_ready' not in st.session_state: st.session_state.model_ready = False

# --- 2. API SETUP ---
def init_gemini(manual_key=None):
    api_key = manual_key
    if not api_key:
        try:
            api_key = st.secrets.get("GOOGLE_API_KEY")
        except Exception: pass

    if api_key:
        try:
            genai.configure(api_key=api_key)
            st.session_state.model_ready = True
            return True
        except Exception: pass
    return False

# --- 3. AI LOGIC ---
def fetch_ai_json(prompt):
    try:
        model = genai.GenerativeModel('gemini-1.5-flash')
        response = model.generate_content(prompt, generation_config={'response_mime_type': 'application/json'})
        return json.loads(response.text)
    except Exception as e:
        st.error(f"AI Error: {e}")
        return None

# --- 4. PPTX EXPORT ---
def create_worksheet_pptx(title, method, comp_qs, graph_io, lesson_qs):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(8.27), Inches(11.69)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    margin = Inches(0.5)
    width = prs.slide_width - (2 * margin)
    
    # Title
    t_box = slide.shapes.add_textbox(margin, Inches(0.2), width, Inches(0.5))
    t_box.text_frame.text = title
    t_box.text_frame.paragraphs[0].font.bold = True
    
    # Method
    m_box = slide.shapes.add_textbox(margin, Inches(0.8), width, Inches(2.2))
    m_box.text_frame.word_wrap = True
    m_box.text_frame.text = f"Method:\n{method}"
    m_box.line.color.rgb = (0,0,0)
    
    # Graph
    graph_io.seek(0)
    slide.shapes.add_picture(graph_io, margin, Inches(3.2), width=width)
    
    # Questions
    q_box = slide.shapes.add_textbox(margin, Inches(7.5), width, Inches(3.8))
    q_box.text_frame.word_wrap = True
    q_box.text_frame.text = f"Questions:\n{chr(10).join(comp_qs)}\n\nAnalysis:\n{lesson_qs}"
    
    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    return ppt_io.getvalue()

# --- 5. MAIN UI ---
st.title("፡ Unified Science Experiment Lab")

with st.sidebar:
    st.header("⚙️ Setup")
    key = st.text_input("Google API Key", type="password")
    if init_gemini(key): st.success("✅ Connected")
    else: st.warning("⚠️ Key Required")
    
    st.divider()
    lo = st.text_area("Learning Objective", "Investigating photosynthesis rates.")
    if st.button("1. Generate Scenario"):
        if st.session_state.model_ready:
            prompt = f"Generate KS3 science scenario for: {lo}. Return JSON: {{'title':'', 'summary':'', 'x_var':'', 'y_var':'', 'z_var_categorical': {{'name':'', 'categories':['','','']}} }}"
            st.session_state.scenario = fetch_ai_json(prompt)

if st.session_state.scenario:
    st.subheader(st.session_state.scenario['title'])
    st.write(st.session_state.scenario['summary'])

    if st.button("2. Generate Data"):
        cols = f"{st.session_state.scenario['x_var']}, {st.session_state.scenario['y_var']}, {st.session_state.scenario['z_var_categorical']['name']}"
        res = fetch_ai_json(f"Generate CSV data for {st.session_state.scenario['title']}. Columns: {cols}. Return JSON: {{'csv_data': ''}}")
        if res: st.session_state.current_df = pd.read_csv(io.StringIO(res['csv_data']))

if st.session_state.current_df is not None:
    fig, ax = plt.subplots(figsize=(10,5))
    sns.scatterplot(data=st.session_state.current_df, x=st.session_state.current_df.columns[0], y=st.session_state.current_df.columns[1], hue=st.session_state.current_df.columns[2], ax=ax)
    ax.grid(True, linestyle=':', alpha=0.6)
    st.pyplot(fig)

    if st.button("3. Finalize Worksheet"):
        comp = fetch_ai_json("Generate method and 5 questions. JSON: {'method_description': '...', 'comprehension_questions': []}")
        lesson = fetch_ai_json("Generate 4 analysis questions. JSON: {'questions_markdown': '...'}")
        if comp and lesson:
            buf = io.BytesIO()
            fig.savefig(buf, format='png', dpi=300)
            ppt = create_worksheet_pptx(st.session_state.scenario['title'], comp['method_description'], comp['comprehension_questions'], buf, lesson['questions_markdown'])
            st.download_button("📥 Download PPTX", ppt, file_name="worksheet.pptx")
